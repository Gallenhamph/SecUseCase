# export.py
import io
import re
import textwrap
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

# --- TEXT CLEANER (CONSOLIDATED) ---
def clean_text(text, mode="pdf"):
    """A master text cleaner that adapts markdown parsing based on the export target."""
    if not text: return ""
    text = text.replace('\xa0', ' ').replace('\t', ' ')
    text = text.replace('“', '"').replace('”', '"').replace('‘', "'").replace('’', "'")
    text = text.replace('–', '-').replace('—', '-')
    text = text.replace('### ', '').replace('## ', '').replace('# ', '')
    
    if mode == "mdr":
        # Strip all markdown links to plain Text (URL) to prevent FPDF monospaced crashing
        text = re.sub(r'\[([^\]]+)\]\((https?://[^\)]+)\)', r'\1 (\2)', text)
        text = text.replace('**', '').replace('*', '').replace('`', '')
    elif mode == "pptx":
        # Convert links to Text: URL for presentation slides
        text = re.sub(r'\[([^\]]+)\]\((https?://[^\)]+)\)', r'\1: \2', text)
        text = text.replace('**', '').replace('*', '')
    
    # "pdf" mode leaves ** and [text](url) intact for native FPDF markdown rendering
    return text.encode('ascii', 'ignore').decode('ascii').strip()

# --- PDF ENGINE ---
class ReportPDF(FPDF):
    def header(self):
        self.set_fill_color(0, 32, 96) 
        self.rect(0, 0, 210, 20, 'F')   
        self.set_y(6)
        self.set_font('helvetica', 'B', 12)
        self.set_text_color(255, 255, 255)
        self.cell(0, 10, 'Threat Modeling & MDR Assessment', align='R')
        self.set_y(25)
        self.set_text_color(0, 0, 0)

    def footer(self):
        self.set_y(-15)
        self.set_font('helvetica', 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f'Page {self.page_no()}', align='C')

def draw_section_header(pdf, title):
    pdf.ln(5)
    pdf.set_font("helvetica", "B", 14)
    pdf.set_text_color(0, 32, 96) 
    pdf.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
    pdf.set_draw_color(200, 200, 200) 
    pdf.set_line_width(0.5)
    pdf.line(pdf.get_x(), pdf.get_y(), 210 - 15, pdf.get_y()) 
    pdf.ln(4)
    pdf.set_text_color(0, 0, 0)

def robust_multi_cell(pdf, w, h, txt, align="L", fill=False):
    try:
        pdf.multi_cell(w=w, h=h, txt=txt, align=align, markdown=True, fill=fill)
    except Exception:
        safe_txt = re.sub(r'\[([^\]]+)\]\((https?://[^\)]+)\)', r'\1 (\2)', txt).replace('**', '').replace('*', '')
        wrap_width = 85 if w == 0 else int(w / 2.0) 
        lines = textwrap.wrap(safe_txt, width=wrap_width, break_long_words=True)
        for line in lines:
            pdf.cell(w=w, h=h, txt=line, align=align, fill=fill, new_x="LMARGIN", new_y="NEXT")

def draw_visual_timeline(pdf, timeline_text):
    if not timeline_text: return
    draw_section_header(pdf, "Attack Timeline & Early MDR Intervention")
    entries = timeline_text.strip().split('\n')
    
    x_node, x_text = 20, 30
    for i, entry in enumerate(entries):
        if not entry.strip() or '|' not in entry: continue
        timestamp, event = entry.split('|', 1)
        if pdf.get_y() > 250: pdf.add_page()
            
        start_y = pdf.get_y()
        pdf.set_fill_color(0, 32, 96)
        pdf.ellipse(x=x_node - 2, y=start_y + 1, w=4, h=4, style='F')
        
        pdf.set_x(x_text)
        pdf.set_font("helvetica", "B", 10)
        pdf.set_text_color(0, 32, 96)
        pdf.cell(w=0, h=6, txt=clean_text(timestamp.strip(), "pdf"), new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("helvetica", "", 10)
        pdf.set_x(x_text)
        
        usable_width = 210 - x_text - 15 
        robust_multi_cell(pdf, usable_width, 5, clean_text(event.strip(), "pdf"))
            
        end_y = pdf.get_y()
        if i < len(entries) - 1:
            pdf.set_draw_color(200, 200, 200)
            pdf.set_line_width(0.5)
            pdf.line(x_node, start_y + 6, x_node, end_y + 2)
        pdf.ln(5)

def create_pdf(inputs, scenario, recs, mdr_case):
    timeline_match = re.search(r'\[TIMELINE_START\](.*?)\[TIMELINE_END\]', scenario, re.DOTALL)
    if timeline_match:
        timeline_text = timeline_match.group(1).strip()
        main_scenario = re.sub(r'\[TIMELINE_START\].*?\[TIMELINE_END\]', '', scenario, flags=re.DOTALL).strip()
    else:
        timeline_text = None
        main_scenario = scenario

    pdf = ReportPDF()
    pdf.add_page()
    
    pdf.set_font("helvetica", "B", 18)
    pdf.cell(w=0, h=12, txt="Cybersecurity Threat & Advisory Report", new_x="LMARGIN", new_y="NEXT", align="C")
    
    pdf.set_font("helvetica", "I", 11)
    pdf.set_text_color(100, 100, 100) 
    pdf.cell(w=0, h=6, txt=f"Prepared for: {inputs['customer_name']}", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.cell(w=0, h=6, txt=f"Presented by: {inputs['consultant_name']}", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_text_color(0, 0, 0) 
    pdf.ln(8)
    
    draw_section_header(pdf, "Client Estate Summary")
    pdf.set_fill_color(245, 245, 245) 
    pdf.set_font("helvetica", "", 10)
    
    summary_text = (
        f"Industry: {inputs['industry']}   |   Users: {inputs['users']}   |   Endpoints: {inputs['endpoints']}\n"
        f"Critical Infrastructure: {inputs['critical_infra']}\n"
        f"M365 License: {inputs['m365_license']}   |   Cloud: {inputs['cloud_env']}\n"
        f"Endpoint: {inputs['endpoint']}   |   Email: {inputs['email']}\n"
        f"Perimeter: {inputs['firewall']} Firewall   |   Identity: {inputs['identity']}\n"
        f"Internal Security: {inputs['in_house_team']}"
    )
    for line in summary_text.split('\n'):
        pdf.cell(w=0, h=7, txt=f"  {line}", new_x="LMARGIN", new_y="NEXT", fill=True)
    pdf.ln(6)
    
    draw_section_header(pdf, "Targeted Threat Narrative & Solutions")
    pdf.set_font("helvetica", "", 10)
    for paragraph in clean_text(main_scenario, "pdf").split('\n'):
        if paragraph.strip():
            robust_multi_cell(pdf, 0, 6, paragraph)
            pdf.ln(2) 
        else:
            pdf.ln(2)
    pdf.ln(6)
    
    if timeline_text:
        draw_visual_timeline(pdf, timeline_text)
        pdf.ln(6)
    
    pdf.add_page() 
    draw_section_header(pdf, "Simulated Sophos MDR Case Log")
    pdf.set_font("courier", "", 9)
    pdf.set_fill_color(240, 248, 255) 
    
    clean_mdr = clean_text(mdr_case, "mdr")
    for line in clean_mdr.split('\n'):
        wrapped_lines = textwrap.wrap(line, width=95, break_long_words=True)
        if not wrapped_lines:
            pdf.cell(w=0, h=5, txt="", new_x="LMARGIN", new_y="NEXT", fill=True)
        for w_line in wrapped_lines:
            pdf.cell(w=0, h=5, txt=f" {w_line}", align="L", fill=True, new_x="LMARGIN", new_y="NEXT")
            
    pdf.ln(6)
    draw_section_header(pdf, "Recommended Security Testing & Advisory")
    pdf.set_font("helvetica", "", 10)
    for r in recs:
        if r.startswith("🛡️") or r.startswith("⚙️"):
            pdf.ln(3)
            robust_multi_cell(pdf, 0, 6, clean_text(r, "pdf"))
        else:
            pdf.set_x(15)
            robust_multi_cell(pdf, 0, 6, clean_text(r, "pdf"))
            pdf.ln(2)
        
    return bytes(pdf.output())

# --- PPTX ENGINE ---
def create_pptx(inputs, scenario, recs, mdr_case):
    DARK_BLUE = RGBColor(0, 32, 96)
    
    timeline_match = re.search(r'\[TIMELINE_START\](.*?)\[TIMELINE_END\]', scenario, re.DOTALL)
    if timeline_match:
        timeline_text = timeline_match.group(1).strip()
        main_scenario = re.sub(r'\[TIMELINE_START\].*?\[TIMELINE_END\]', '', scenario, flags=re.DOTALL).strip()
    else:
        timeline_text = None
        main_scenario = scenario

    prs = Presentation()
    
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Threat Modeling & MDR Assessment"
    slide1.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    slide1.placeholders[1].text = f"Prepared for: {inputs['customer_name']}\nPresented by: {inputs['consultant_name']}\n{inputs['industry']} Sector"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = "Client Estate Overview"
    slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    tf2 = slide2.shapes.placeholders[1].text_frame
    tf2.clear()
    
    details = [
        f"Target Industry: {inputs['industry']}",
        f"Total User Base: {inputs['users']} Users (Security Savviness: {inputs['savviness']})",
        f"Infrastructure: {inputs['endpoints']} Endpoints | {inputs['servers']} Servers",
        f"Crown Jewels: {inputs['critical_infra']}",
        f"Internal Security Team: {inputs['in_house_team']}",
        "Current Technology Stack:",
        f"  • Endpoint: {inputs['endpoint']}",
        f"  • Perimeter: {inputs['firewall']}",
        f"  • Identity: {inputs['identity']}",
        f"  • Email: {inputs['email']}",
        f"  • M365/Cloud: {inputs['m365_license']} | {inputs['cloud_env']}"
    ]
    for d in details:
        p = tf2.add_paragraph()
        p.text = d
        p.font.size = Pt(14)
        if "Technology Stack:" in d:
            p.font.bold = True
            p.space_before = Pt(14)
    
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    slide3.shapes.title.text = "Threat Narrative: Executive Summary"
    slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    tf3 = slide3.shapes.placeholders[1].text_frame
    tf3.word_wrap = True 
    tf3.clear()
    
    paras = [p for p in main_scenario.split('\n') if p.strip()]
    exec_summary = paras[:2] if len(paras) >= 2 else paras
    for para in exec_summary:
        p = tf3.add_paragraph()
        p.text = clean_text(para, "pptx")
        p.font.size = Pt(14) 
        p.space_after = Pt(14)
    
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Attack Timeline & MDR Intervention"
    slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    if timeline_text:
        entries = timeline_text.strip().split('\n')
        x_node, x_text, y_offset, text_width = Inches(0.8), Inches(1.3), Inches(1.8), Inches(8.0)

        for i, entry in enumerate(entries):
            if '|' not in entry: continue
            timestamp, event = entry.split('|', 1)

            node = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x_node, y_offset, Inches(0.2), Inches(0.2))
            node.fill.solid()
            node.fill.fore_color.rgb = DARK_BLUE
            node.line.color.rgb = DARK_BLUE

            txBox = slide4.shapes.add_textbox(x_text, y_offset - Inches(0.1), text_width, Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            p1.text = clean_text(timestamp.strip(), "pptx")
            p1.font.bold = True
            p1.font.color.rgb = DARK_BLUE
            p1.font.size = Pt(12)

            p2 = tf.add_paragraph()
            p2.text = clean_text(event.strip(), "pptx")
            p2.font.size = Pt(12)

            line_height = Inches(0.8)
            if len(event) > 80: line_height = Inches(1.2)
            if len(event) > 150: line_height = Inches(1.5)
            
            next_y_offset = y_offset + line_height

            if i < len(entries) - 1:
                connector = slide4.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_node + Inches(0.1), y_offset + Inches(0.2), x_node + Inches(0.1), next_y_offset)
                connector.line.color.rgb = RGBColor(150, 150, 150)
            y_offset = next_y_offset

    slide5 = prs.slides.add_slide(bullet_slide_layout)