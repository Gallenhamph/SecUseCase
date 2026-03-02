# app.py
import streamlit as st
import io
import re
import random
import textwrap
from google import genai
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Pt
from prompts import SYSTEM_PERSONA, build_scenario_prompt, build_mdr_case_prompt

# --- BACKEND LOGIC ---
class CyberScenarioGenerator:
    def __init__(self, api_key):
        self.api_key = api_key
        if self.api_key and self.api_key != "YOUR_API_KEY_HERE":
            self.client = genai.Client(api_key=self.api_key)
        else:
            self.client = None
    
    def fetch_osint(self, vendor):
        simulated_osint = {
            "Fortinet": "Active exploitation of FortiOS SSL-VPN vulnerabilities (e.g., CVE-2023-27997, CVE-2024-21762) by state-sponsored actors to deploy custom implants and bypass pre-authentication filters.",
            "Palo Alto": "Rising trend of threat actors exploiting unpatched PAN-OS GlobalProtect interfaces (e.g., CVE-2024-3400) to achieve unauthenticated remote code execution and establish persistent reverse shells.",
            "Cisco": "Exploitation of AnyConnect and IOS XE zero-days (e.g., CVE-2023-20198), leading to privilege escalation and the deployment of malicious Lua-based web shells on edge appliances.",
            "Check Point": "Targeted attacks exploiting Check Point Security Gateway vulnerabilities (e.g., CVE-2024-24919) to extract Active Directory hashes and establish persistent VPN sessions.",
            "SonicWall": "Continued exploitation of SMA 100 series appliances (e.g., CVE-2021-20038) using credential stuffing and unpatched firmware to deploy ransomware directly into the DMZ.",
            "WatchGuard": "Historical targeting by botnets exploiting unpatched privilege escalation flaws (e.g., CVE-2022-26318) to maintain long-term, stealthy persistence on edge devices.",
            "Barracuda": "Sophisticated threat actors exploiting Email Security Gateway (ESG) zero-days (e.g., CVE-2023-2868) to deploy data exfiltration malware and backdoors.",
            "Juniper": "Exploitation of Junos OS J-Web vulnerabilities (e.g., CVE-2023-36844) allowing unauthenticated attackers to upload arbitrary files and execute code as root.",
            "CrowdStrike": "Advanced adversaries increasingly utilizing custom bootloaders and kernel-level drivers (BYOVD - Bring Your Own Vulnerable Driver) to blind Falcon sensors (Reference: Elastic Security Labs BYOVD research).",
            "Microsoft Defender": "High reliance on 'Living off the Land' (LotL) techniques and obfuscated PowerShell scripts to evade standard Defender ASR rules and execute fileless malware.",
            "SentinelOne": "Threat actors utilizing highly obfuscated, fragmented shellcode and direct syscalls to evade SentinelOne's behavioral AI engines.",
            "Trend Micro": "Exploitation of legacy Apex One vulnerabilities (e.g., CVE-2022-40139) and exploitation of exclusion lists to deploy ransomware payloads undetected.",
            "Symantec": "Bypass of legacy signature-based protections using polymorphic malware families and living-off-the-land binaries (LOLBins).",
            "Okta": "Surge in highly sophisticated Adversary-in-the-Middle (AiTM) phishing kits capturing Okta session cookies and bypassing multi-factor authentication entirely (Reference: CISA Advisory AA23-320A).",
            "Microsoft Entra ID (Azure AD)": "Widespread MFA fatigue (push bombing) attacks and illicit consent grants involving malicious OAuth applications to maintain persistent access to Microsoft 365 environments.",
            "Cisco Duo": "Targeting of telephony-based authentication (SMS/Voice) via SIM swapping, alongside localized push-notification fatigue campaigns.",
            "Mimecast": "Massive increase in Quishing (QR Code Phishing) and HTML smuggling campaigns that successfully bypass Mimecast's URL rewriting and attachment sandboxing (Reference: IBM X-Force Quishing Trends).",
            "Proofpoint": "Threat actors leveraging highly customized, evasive PDF documents containing embedded malicious links that bypass TAP (Targeted Attack Protection) analysis.",
            "AWS": "Exploitation of overly permissive IAM roles via SSRF vulnerabilities on public-facing EC2 instances, leading to environment-wide administrative compromise.",
            "Microsoft Azure": "Abuse of Azure Automation Runbooks and extraction of Managed Identity tokens to pivot laterally across the Azure environment.",
            "GCP": "Targeting of exposed service account keys embedded in developer repositories to access Google Cloud Storage buckets and BigQuery datasets."
        }
        return simulated_osint.get(vendor, "")

    def generate_recommendations(self, inputs):
        recs = []
        
        recs.append("🛡️ **SECURITY ASSESSMENTS & ADVISORY**")
        
        if inputs['in_house_team'] == "Yes (24/7)" and inputs['savviness'] in ["Medium", "High"]:
             recs.append("• [Secureworks Adversary Exercises (Red Teaming)](https://www.secureworks.com/services/offensive-security): Emulate a sophisticated adversary to stress-test your mature 24/7 SOC and validate detection capabilities across the kill chain.")
             recs.append("• [Secureworks Threat Hunting Assessment](https://www.secureworks.com/services/threat-hunting): Proactively search your environment for undetected threats or persistence mechanisms that may have bypassed your existing defenses.")
        elif inputs['in_house_team'] != "No":
             recs.append("• [Sophos Internal Penetration Testing](https://www.sophos.com/en-us/services/penetration-testing): Simulate an attacker who has bypassed the perimeter to test domain compromise, internal lateral movement, and existing defenses.")
             recs.append("• [Secureworks Tabletop Exercises](https://www.secureworks.com/services/incident-response-readiness): Ensure leadership and the internal security team are aligned on communication, legal, and operational procedures during a crisis.")
        else:
            recs.append("• [Sophos Emergency Incident Response Retainer](https://www.sophos.com/en-us/services/incident-response-retainer): Crucial for organizations without dedicated internal IR teams to guarantee SLAs and immediate assistance during a live breach.")
            recs.append("• [Secureworks Ransomware Readiness Assessment](https://www.secureworks.com/services/incident-response-readiness): Evaluate your organization's preparedness to defend against, endure, and recover from a targeted ransomware event.")

        if inputs['public_web_apps']:
            recs.append("• [Sophos Web Application Security Assessment](https://www.sophos.com/en-us/services/penetration-testing): Identify coding flaws (e.g., SQLi, XSS) in your public-facing web applications before attackers exploit them to access backend databases.")
            recs.append("• [Sophos External Penetration Testing](https://www.sophos.com/en-us/services/penetration-testing): Manual, tester-driven attempts to breach your internet-facing assets, moving beyond automated vulnerability scanning.")

        if inputs['physical_locations'] > 1:
             recs.append("• [Sophos Wireless Network Penetration Testing](https://www.sophos.com/en-us/services/penetration-testing): Evaluate wireless security across your physical locations, testing for rogue access points and weak encryption.")

        if "Cloud" in inputs['cloud_env'] or inputs['cloud_env'] in ["AWS", "Microsoft Azure", "GCP"]:
            recs.append(f"• [Sophos Cloud Security Assessment](https://www.sophos.com/en-us/services/cloud-security-posture-management): Audit your {inputs['cloud_env']} environment for misconfigurations, overly permissive IAM roles, and compliance violations.")

        if inputs['savviness'] == "Low":
            recs.append("• [Sophos Phish Threat (Social Engineering Simulation)](https://www.sophos.com/en-us/products/phish-threat): Conduct targeted phishing and social engineering exercises to baseline and improve employee security awareness.")

        if "Active Directory" in inputs['identity'] or "Entra ID" in inputs['identity']:
             recs.append("• [Sophos Active Directory Security Assessment](https://www.sophos.com/en-us/services/compromise-assessment): Identify architectural weaknesses, excessive privileges, and misconfigurations in your core identity platform.")

        if inputs['servers'] > 50:
            recs.append("• [Sophos Compromise Assessment](https://www.sophos.com/en-us/services/compromise-assessment): Proactively hunt for existing persistence mechanisms or dormant threats currently hiding in your data center.")

        recs.append("\n⚙️ **RECOMMENDED SOPHOS SOLUTIONS**")
        
        if inputs['m365_license'] != "None / On-Prem Only":
            recs.append(f"• [Sophos MDR for Microsoft 365](https://www.sophos.com/en-us/products/mdr): Maximize your {inputs['m365_license']} investment. Sophos ingests telemetry directly from Microsoft Graph Security, Entra ID, and Defender to correlate Microsoft alerts with cross-domain threat intelligence, stopping attacks that bypass native Microsoft controls.")

        recs.append("• [Sophos Managed Risk](https://www.sophos.com/en-us/products/managed-risk): Implement continuous external attack surface management to discover and prioritize exposed vulnerabilities across your evolving tech stack before they can be weaponized.")

        if inputs['identity'] not in ["None / Local Only", "On-Prem Active Directory"]:
            recs.append(f"• [Sophos ITDR (Identity Threat Detection and Response)](https://www.sophos.com/en-us/products/mdr): Integrate telemetry directly from {inputs['identity']} to detect compromised credentials, anomalous logins, and lateral movement tied to user identities before endpoints are even touched.")

        if inputs['firewall'] != "Sophos" or inputs['servers'] > 20:
            recs.append("• [Sophos NDR (Network Detection and Response)](https://www.sophos.com/en-us/products/network-detection-and-response): Analyze network traffic for rogue devices, unprotected assets, and insider threats. This is critical for monitoring lateral movement across the network in a 'Bring Your Own Tech' environment.")

        if inputs['endpoint'] != "Sophos":
             recs.append(f"• [Sophos Intercept X Advanced with XDR](https://www.sophos.com/en-us/products/endpoint-antivirus): Consolidate your endpoint stack by replacing {inputs['endpoint']}. This provides Sophos MDR analysts with native, deep-level remediation capabilities rather than just third-party telemetry.")
             
        if inputs['email'] != "Sophos":
             recs.append(f"• [Sophos Email Security](https://www.sophos.com/en-us/products/email-security): Integrate advanced phishing protection and post-delivery remediation natively into your MDR ecosystem, reducing the risk of social engineering attacks that bypass {inputs['email']}.")

        return recs

    def call_llm(self, prompt):
        if not self.client:
            return "⚠️ Error: Please enter a valid Gemini API Key in the application code or via Streamlit secrets."
            
        full_prompt = f"{SYSTEM_PERSONA}\n\n{prompt}"
        
        try:
            response = self.client.models.generate_content(
                model='gemini-2.5-flash',
                contents=full_prompt
            )
            return response.text
        except Exception as e:
            return f"⚠️ An error occurred while communicating with the Gemini API: {e}"

# --- EXPORT LOGIC ---
class ReportPDF(FPDF):
    def header(self):
        self.set_fill_color(0, 32, 96) 
        self.rect(0, 0, 210, 20, 'F')   
        self.set_y(6)
        self.set_font('helvetica', 'B', 12)
        self.set_text_color(255, 255, 255)
        self.cell(0, 10, 'Threat Modeling & MDR Assessment', align='R')
        self.set_y(25)
        self.set_text_color(0, 0, 0)

    def footer(self):
        self.set_y(-15)
        self.set_font('helvetica', 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f'Page {self.page_no()}', align='C')

def draw_section_header(pdf, title):
    pdf.ln(5)
    pdf.set_font("helvetica", "B", 14)
    pdf.set_text_color(0, 32, 96) 
    pdf.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
    pdf.set_draw_color(200, 200, 200) 
    pdf.set_line_width(0.5)
    pdf.line(pdf.get_x(), pdf.get_y(), 210 - 15, pdf.get_y()) 
    pdf.ln(4)
    pdf.set_text_color(0, 0, 0)

def clean_text(text):
    if not text: return ""
    text = text.replace('\xa0', ' ').replace('\t', ' ')
    
    text = re.sub(r'\[([^\]]+)\]\((https?://[^\)]+)\)', r'\1 (\2)', text)
    
    text = text.replace('**', '').replace('*', '').replace('#', '')
    text = text.replace('“', '"').replace('”', '"').replace('‘', "'").replace('’', "'")
    text = text.replace('–', '-').replace('—', '-')
    text = text.replace('### ', '').replace('## ', '').replace('# ', '') 
    text = text.encode('ascii', 'ignore').decode('ascii')
    return text.strip()

def clean_pptx_text(text):
    text = clean_text(text)
    text = re.sub(r'\[([^\]]+)\]\((https?://[^\)]+)\)', r'\1: \2', text)
    text = text.replace('**', '').replace('*', '')
    return text.strip()

def robust_multi_cell(pdf, w, h, txt, align="L", fill=False):
    try:
        pdf.multi_cell(w=w, h=h, txt=txt, align=align, markdown=True, fill=fill)
    except Exception:
        safe_txt = re.sub(r'\[([^\]]+)\]\((https?://[^\)]+)\)', r'\1', txt).replace('**', '').replace('*', '')
        wrap_width = 90 if w == 0 else int(w / 1.8) 
        lines = textwrap.wrap(safe_txt, width=wrap_width)
        for line in lines:
            pdf.cell(w=w, h=h, txt=line, align=align, fill=fill, new_x="LMARGIN", new_y="NEXT")

def write_safe_text(pdf, text, font_family="helvetica"):
    pdf.set_font(font_family, "", 10)
    paragraphs = text.split('\n')
    for paragraph in paragraphs:
        if paragraph.strip():
            robust_multi_cell(pdf, 0, 6, paragraph)
            pdf.ln(2) 
        else:
            pdf.ln(2)

def draw_visual_timeline(pdf, timeline_text):
    if not timeline_text: return
    
    draw_section_header(pdf, "Attack Timeline & Early MDR Intervention")
    entries = timeline_text.strip().split('\n')
    
    x_node = 20
    x_text = 30
    
    for i, entry in enumerate(entries):
        if not entry.strip() or '|' not in entry: 
            continue
            
        timestamp, event = entry.split('|', 1)
        
        if pdf.get_y() > 250:
            pdf.add_page()
            
        start_y = pdf.get_y()
        
        pdf.set_fill_color(0, 32, 96)
        pdf.ellipse(x=x_node - 2, y=start_y + 1, w=4, h=4, style='F')
        
        pdf.set_x(x_text)
        pdf.set_font("helvetica", "B", 10)
        pdf.set_text_color(0, 32, 96)
        pdf.cell(w=0, h=6, txt=clean_text(timestamp.strip()), new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("helvetica", "", 10)
        pdf.set_x(x_text)
        
        usable_width = 210 - x_text - 15 
        robust_multi_cell(pdf, usable_width, 5, clean_text(event.strip()))
            
        end_y = pdf.get_y()
        
        if i < len(entries) - 1:
            pdf.set_draw_color(200, 200, 200)
            pdf.set_line_width(0.5)
            pdf.line(x_node, start_y + 6, x_node, end_y + 2)
            
        pdf.ln(5)

def create_pdf(inputs, scenario, recs, mdr_case):
    timeline_match = re.search(r'\[TIMELINE_START\](.*?)\[TIMELINE_END\]', scenario, re.DOTALL)
    if timeline_match:
        timeline_text = timeline_match.group(1).strip()
        main_scenario = re.sub(r'\[TIMELINE_START\].*?\[TIMELINE_END\]', '', scenario, flags=re.DOTALL).strip()
    else:
        timeline_text = None
        main_scenario = scenario

    pdf = ReportPDF()
    pdf.add_page()
    
    pdf.set_font("helvetica", "B", 18)
    pdf.cell(w=0, h=12, txt="Cybersecurity Threat & Advisory Report", new_x="LMARGIN", new_y="NEXT", align="C")
    
    pdf.set_font("helvetica", "I", 11)
    pdf.set_text_color(100, 100, 100) 
    pdf.cell(w=0, h=6, txt=f"Prepared for: {inputs['customer_name']}", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.cell(w=0, h=6, txt=f"Presented by: {inputs['consultant_name']}", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_text_color(0, 0, 0) 
    pdf.ln(8)
    
    draw_section_header(pdf, "Client Estate Summary")
    pdf.set_fill_color(245, 245, 245) 
    pdf.set_font("helvetica", "", 10)
    
    summary_text = (
        f"Industry: {inputs['industry']}   |   Users: {inputs['users']}   |   Endpoints: {inputs['endpoints']}\n"
        f"Critical Infrastructure: {inputs['critical_infra']}\n"
        f"M365 License: {inputs['m365_license']}   |   Cloud: {inputs['cloud_env']}\n"
        f"Endpoint: {inputs['endpoint']}   |   Email: {inputs['email']}\n"
        f"Perimeter: {inputs['firewall']} Firewall   |   Identity: {inputs['identity']}\n"
        f"Internal Security: {inputs['in_house_team']}"
    )
    
    for line in summary_text.split('\n'):
        pdf.cell(w=0, h=7, txt=f"  {line}", new_x="LMARGIN", new_y="NEXT", fill=True)
    
    pdf.ln(6)
    
    draw_section_header(pdf, "Targeted Threat Narrative & Solutions")
    write_safe_text(pdf, clean_text(main_scenario))
    pdf.ln(6)
    
    if timeline_text:
        draw_visual_timeline(pdf, timeline_text)
        pdf.ln(6)
    
    pdf.add_page() 
    draw_section_header(pdf, "Simulated Sophos MDR Case Log")
    
    pdf.set_font("courier", "", 9)
    pdf.set_fill_color(240, 248, 255) 
    
    clean_mdr = clean_text(mdr_case)
    for line in clean_mdr.split('\n'):
        robust_multi_cell(pdf, 0, 5, f" {line}", fill=True)
            
    pdf.ln(6)
    
    draw_section_header(pdf, "Recommended Security Testing & Advisory")
    for r in recs:
        if r.startswith("🛡️") or r.startswith("⚙️"):
            pdf.ln(3)
            write_safe_text(pdf, clean_text(r))
        else:
            pdf.set_x(15)
            robust_multi_cell(pdf, 0, 6, clean_text(r))
            pdf.ln(2)
        
    return bytes(pdf.output())

def create_pptx(inputs, scenario, recs, mdr_case):
    scenario_clean = re.sub(r'\[TIMELINE_START\]', '\n--- Attack Timeline ---\n', scenario)
    scenario_clean = re.sub(r'\[TIMELINE_END\]', '', scenario_clean)

    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Threat Modeling & MDR Assessment"
    subtitle.text = f"Prepared for: {inputs['customer_name']}\nPresented by: {inputs['consultant_name']}\n{inputs['industry']} Sector"
    
    paragraphs = [p for p in scenario_clean.split('\n') if p.strip()]
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = "Attack Narrative - Phase 1"
    tf2 = slide2.shapes.placeholders[1].text_frame
    tf2.word_wrap = True 
    tf2.clear()
    
    for para in paragraphs[:2]:
        p = tf2.add_paragraph()
        p.text = clean_pptx_text(para)
        p.font.size = Pt(14) 
        p.space_after = Pt(10)
        
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    slide3.shapes.title.text = "Attack Narrative - Phase 2"
    tf3 = slide3.shapes.placeholders[1].text_frame
    tf3.word_wrap = True 
    tf3.clear()
    
    for para in paragraphs[2:5]:
        p = tf3.add_paragraph()
        p.text = clean_pptx_text(para)
        p.font.size = Pt(14) 
        p.space_after = Pt(10)
    
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    slide4.shapes.title.text = "Simulated MDR Investigation"
    tf4 = slide4.shapes.placeholders[1].text_frame
    tf4.word_wrap = True
    
    analysis_match = re.search(r'//Analysis:(.*?)//Response Actions:', mdr_case, re.DOTALL)
    if analysis_match:
        analysis_text = clean_pptx_text(analysis_match.group(1).strip())
    else:
        analysis_text = clean_pptx_text(mdr_case[:500]) + "..."
        
    p4 = tf4.paragraphs[0]
    p4.text = analysis_text
    p4.font.size = Pt(14)
    
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    slide5.shapes.title.text = "Testing & Advisory Recommendations"
    tf5 = slide5.shapes.placeholders[1].text_frame
    tf5.word_wrap = True
    tf5.clear() 
    
    for r in recs:
        p5 = tf5.add_paragraph()
        if r.startswith("🛡️") or r.startswith("⚙️"):
            p5.text = clean_pptx_text(r)
            p5.font.bold = True
            p5.font.size = Pt(16)
            p5.level = 0
            p5.space_before = Pt(10)
        else:
            p5.text = clean_pptx_text(r).lstrip('•').strip()
            p5.font.size = Pt(12)
            p5.level = 1
            
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    return pptx_stream.getvalue()

# --- STREAMLIT FRONTEND ---
st.set_page_config(page_title="MDR & Testing Scenario Generator", page_icon="🛡️", layout="wide")

st.title("🛡️ MDR & Offensive Security Scenario Generator")
st.markdown("Generate highly tailored cyberattack scenarios, complete with mock Sophos Central case logs.")

try:
    gemini_key = st.secrets["GEMINI_API_KEY"]
except KeyError:
    st.error("Missing API Key! Please add GEMINI_API_KEY to your .streamlit/secrets.toml file.")
    gemini_key = None

app_engine = CyberScenarioGenerator(api_key=gemini_key)

with st.sidebar:
    st.header("📋 Engagement Details")
    customer_name = st.text_input("Customer Name", "Acme Corp")
    consultant_name = st.text_input("Consultant Name", "Jane Doe")
    
    st.header("🏢 Client Estate Details")
    industry = st.selectbox("Industry Vertical", ["Healthcare", "Finance", "Manufacturing", "Retail", "Technology", "Education"])
    critical_infra = st.text_input("Critical Infrastructure/Crown Jewels", "Patient Records Database")
    
    st.subheader("👥 The Human Element")
    users = st.number_input("Number of Users", min_value=1, value=500)
    savviness = st.select_slider("User Security Savviness", options=["Low", "Medium", "High"])
    in_house_team = st.radio("In-House Security Team?", ["No", "Yes (9-to-5)", "Yes (24/7)"])
    
    st.subheader("💻 Technology Stack")
    endpoints = st.number_input("Number of Endpoints", min_value=1, value=600)
    servers = st.number_input("Number of Servers", min_value=1, value=50)
    
    endpoint = st.selectbox("Endpoint Security", ["Sophos", "Microsoft Defender", "CrowdStrike", "SentinelOne", "Trend Micro", "Symantec", "Trellix", "Blackberry Cylance", "Other"])
    firewall = st.selectbox("Firewall Vendor", ["Fortinet", "Palo Alto", "Cisco", "Check Point", "SonicWall", "WatchGuard", "Juniper", "Barracuda", "Forcepoint", "Sophos", "Other"])
    identity = st.selectbox("Identity Provider", ["Microsoft Entra ID (Azure AD)", "Okta", "Cisco Duo", "Ping Identity", "On-Prem Active Directory", "None / Local Only"])
    m365_license = st.selectbox("Microsoft 365 Licensing", ["None / On-Prem Only", "M365 Business Basic/Standard", "M365 Business Premium", "Office 365 E3 / M365 E3", "Microsoft 365 E5"])
    email = st.selectbox("Email Security", ["Sophos", "Microsoft Defender for Office 365", "Mimecast", "Proofpoint", "Barracuda", "Other"])
    cloud_env = st.selectbox("Cloud Infrastructure", ["AWS", "Microsoft Azure", "GCP", "Multi-Cloud", "None (Fully On-Prem)"])

    st.subheader("🎯 Additional Attack Surfaces")
    physical_locations = st.number_input("Number of Physical Offices/Locations", min_value=1, value=3)
    public_web_apps = st.checkbox("Host Public-Facing Web Applications?")
    
    st.subheader("🛠️ Custom Scenario Override")
    custom_scenario = st.text_area("Specific Threat/Use Case (Optional)", placeholder="e.g., 'How would Sophos MDR detect a BlackBasta ransomware deployment via a compromised VPN?' Leave blank for a random scenario.")

    generate_btn = st.button("Generate Full Scenario", type="primary")

# --- SESSION STATE GENERATION ---
if generate_btn:
    client_inputs = {
        "customer_name": customer_name, "consultant_name": consultant_name,
        "industry": industry, "users": users, "savviness": savviness, 
        "endpoints": endpoints, "servers": servers, "critical_infra": critical_infra,
        "endpoint": endpoint, "firewall": firewall, "identity": identity, "m365_license": m365_license, "email": email, "cloud_env": cloud_env,
        "in_house_team": in_house_team, "physical_locations": physical_locations, "public_web_apps": public_web_apps
    }
    
    # Store dynamic inputs in session state so file names match the generated output
    st.session_state['client_inputs'] = client_inputs
    st.session_state['customer_name'] = customer_name
    
    attack_vectors = [
        "Highly targeted spear-phishing campaign using a malicious PDF attachment (T1566.001)",
        "Adversary-in-the-Middle (AiTM) proxy attack defeating standard MFA via a fake login page (T1556)",
        "Voice Phishing (Vishing) the IT Helpdesk to fraudulently reset a user's MFA device (T1566.004)",
        "Social engineering via LinkedIn/Slack delivering a malicious payload disguised as a resume (T1566.003)",
        "Exploitation of a zero-day vulnerability in a public-facing web application (T1190)",
        "Exploitation of an unpatched, legacy VPN appliance leading to internal access (T1133)",
        "Password spraying attack against legacy authentication protocols lacking MFA enforcement (T1110.003)",
        "Default credentials left active on an internet-facing IoT or edge network device (T1078.001)",
        "Compromised third-party IT contractor / Supply Chain Compromise via remote access tools (T1195)",
        "Malicious update pushed through a compromised third-party software vendor (T1195.002)",
        "Lateral pivot into the network originating from a compromised trusted vendor's environment (T1199)",
        "Compromised Cloud Infrastructure via hardcoded API keys accidentally leaked on GitHub (T1078.004)",
        "Session hijacking via stolen browser cookies purchased on the dark web, bypassing MFA entirely (T1539)",
        "Malicious insider abusing legitimate administrative privileges to disable security tooling (T1078.003)",
        "Physical 'USB Drop' attack in the company parking lot leading to a reverse shell beacon (T1200)"
    ]
    
    selected_vector = random.choice(attack_vectors)
    
    with st.spinner("Analyzing estate and generating narrative with Gemini 2.5 Flash..."):
        osint_list = [
            app_engine.fetch_osint(endpoint),
            app_engine.fetch_osint(firewall),
            app_engine.fetch_osint(identity),
            app_engine.fetch_osint(email),
            app_engine.fetch_osint(cloud_env)
        ]
        osint_data = " ".join([x for x in osint_list if x])
        
        narrative_prompt = build_scenario_prompt(client_inputs, osint_data, selected_vector, custom_scenario)
        scenario = app_engine.call_llm(narrative_prompt)
        
        case_prompt = build_mdr_case_prompt(client_inputs, scenario)
        mdr_case = app_engine.call_llm(case_prompt)
        
        recs = app_engine.generate_recommendations(client_inputs)
        
        # Save results to session state
        st.session_state['scenario'] = scenario
        st.session_state['mdr_case'] = mdr_case
        st.session_state['recs'] = recs
        
        # Pre-generate bytes for persistence
        if "⚠️ Error" not in scenario and "⚠️ An error" not in scenario:
            st.session_state['pdf_bytes'] = create_pdf(client_inputs, scenario, recs, mdr_case)
            st.session_state['pptx_bytes'] = create_pptx(client_inputs, scenario, recs, mdr_case)
        else:
            st.session_state['pdf_bytes'] = None
            st.session_state['pptx_bytes'] = None
            
    st.success("Analysis Complete!")

# --- UI RENDERING FROM SESSION STATE ---
if 'scenario' in st.session_state:
    # Retrieve from cache
    scenario = st.session_state['scenario']
    mdr_case = st.session_state['mdr_case']
    recs = st.session_state['recs']
    cached_customer_name = st.session_state['customer_name']
    
    ui_scenario = re.sub(r'\[TIMELINE_START\]', '\n#### Attack Timeline\n', scenario)
    ui_scenario = re.sub(r'\[TIMELINE_END\]', '', ui_scenario)
    
    tab1, tab2, tab3 = st.tabs(["📝 Threat Narrative", "🛡️ Sophos Central MDR Log", "🎯 Recommendations"])
    
    with tab1:
        st.subheader(f"Threat Narrative & Solutions for {cached_customer_name}")
        st.write(ui_scenario)
        
    with tab2:
        st.subheader("Simulated MDR Investigation")
        st.info("This output mimics the Case Details view a customer would receive in Sophos Central.")
        st.write(mdr_case)
        
    with tab3:
        st.subheader("Recommended Security Testing & Advisory")
        for rec in recs:
            if rec.startswith("🛡️") or rec.startswith("⚙️"):
                st.markdown(f"#### {rec}")
            else:
                st.markdown(rec)
        
    st.divider()
    
    if st.session_state.get('pdf_bytes') and st.session_state.get('pptx_bytes'):
        st.subheader("📥 Export Client Deliverables")
        dl_col1, dl_col2 = st.columns(2)
        
        with dl_col1:
            st.download_button(
                label="📄 Download PDF Report",
                data=st.session_state['pdf_bytes'],
                file_name=f"{cached_customer_name.replace(' ', '_')}_MDR_Report.pdf",
                mime="application/pdf"
            )
            
        with dl_col2:
            st.download_button(
                label="📊 Download PowerPoint Deck",
                data=st.session_state['pptx_bytes'],
                file_name=f"{cached_customer_name.replace(' ', '_')}_MDR_Deck.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )